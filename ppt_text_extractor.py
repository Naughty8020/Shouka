from pptx import Presentation

def extract_text(ppt_path):   # ← 関数名を変更
    prs = Presentation(ppt_path)
    all_text = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        all_text.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))

    return "\n\n".join(all_text)
