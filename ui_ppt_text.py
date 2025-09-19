# ui_app_with_slide_view.py
import sys, os, subprocess
from PySide6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QPushButton, QTextEdit, QLabel, QFileDialog, QComboBox, QScrollArea
)
from PySide6.QtGui import QPixmap
from pptx import Presentation
from ppt_text_extractor import extract_text
from pptx.util import Inches

class PPTApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Shouka PPT AI Helper")
        self.setGeometry(200, 200, 1200, 700)

        self.ppt_path = ""
        self.edited_ppt_path = None
        self.slides_pix = []
        self.slide_index = 0

        main_widget = QWidget()
        self.setCentralWidget(main_widget)
        layout = QHBoxLayout(main_widget)

        # ---------------- 左レイアウト ----------------
        left_layout = QVBoxLayout()
        self.open_btn = QPushButton("📂資料をSHOUKAさせる")
        self.open_btn.clicked.connect(self.load_ppt)
        self.ppt_path_label = QLabel("資料のパス: なし")
        self.input_text = QTextEdit()

        self.save_btn = QPushButton("💾SHOUKAした資料を出力・作成")
        self.save_btn.clicked.connect(self.save_ppt)

        self.open_in_app_btn = QPushButton("💻資料をで開く・確認")
        self.open_in_app_btn.clicked.connect(self.open_in_powerpoint)

        # 表示切替用
        self.view_box = QComboBox()
        self.view_box.addItems(["元のPPT", "編集後PPT"])
        self.view_box.currentIndexChanged.connect(self.update_slide_view)

        # スライド表示
        self.slide_label = QLabel()
        self.slide_label.setFixedHeight(300)
        self.scroll = QScrollArea()
        self.scroll.setWidgetResizable(True)
        self.scroll.setWidget(self.slide_label)

        left_layout.addWidget(self.open_btn)
        left_layout.addWidget(self.ppt_path_label)
        left_layout.addWidget(QLabel("テキスト"))
        left_layout.addWidget(self.input_text)
        left_layout.addWidget(self.save_btn)
        left_layout.addWidget(self.open_in_app_btn)
        left_layout.addWidget(QLabel("表示PPT選択"))
        left_layout.addWidget(self.view_box)
   

        # ---------------- 中央レイアウト ----------------
        center_layout = QVBoxLayout()
        self.mode_box = QComboBox()
        self.mode_box.addItems(["翻訳", "トーンアップ"])
        self.engine_box = QComboBox()
        self.engine_box.addItems(["クラウドAI", "ローカルAI"])
        self.run_btn = QPushButton("🚀 実行")
        self.run_btn.clicked.connect(self.run_ai)
        center_layout.addWidget(QLabel("モード選択"))
        center_layout.addWidget(self.mode_box)
        center_layout.addWidget(QLabel("AI エンジン"))
        center_layout.addWidget(self.engine_box)
        center_layout.addWidget(self.run_btn)

        # ---------------- 右レイアウト ----------------
        right_layout = QVBoxLayout()
        self.output_text = QTextEdit()
        right_layout.addWidget(QLabel("SHOUKA結果"))
        right_layout.addWidget(self.output_text)

        # ---------------- 全体レイアウト ----------------
        layout.addLayout(left_layout, 4)
        layout.addLayout(center_layout, 1)
        layout.addLayout(right_layout, 4)

    # ---------------- PPT読み込み ----------------
    def load_ppt(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "PPTを選択", "", "PowerPoint Files (*.pptx)")
        if file_path:
            self.ppt_path = file_path
            self.ppt_path_label.setText(f"PPTパス: {file_path}")
            text = extract_text(file_path)
            self.input_text.setText(text)
            self.convert_ppt_to_images(self.ppt_path)
            self.show_slide()

    # ---------------- 既定アプリで PPT を開く ----------------
    def open_in_powerpoint(self):
        if not self.view_box.currentText() == "編集後PPT" or not self.edited_ppt_path:
            file_to_open = self.ppt_path
        else:
            file_to_open = self.edited_ppt_path

        if not file_to_open:
            return
        if sys.platform == "darwin":
            subprocess.run(["open", file_to_open])
        elif sys.platform == "win32":
            os.startfile(file_to_open)
        else:
            subprocess.run(["xdg-open", file_to_open])

    # ---------------- AI処理（仮） ----------------
    def run_ai(self):
        mode = self.mode_box.currentText()
        engine = self.engine_box.currentText()
        input_text = self.input_text.toPlainText()
        result = f"[{engine}] {mode} 処理結果：\n\n" + input_text.upper()
        self.output_text.setText(result)

    # ---------------- フォント保持でテキスト置き換え ----------------
    def replace_text_keep_font(self, shape, new_text):
        if not shape.has_text_frame:
            return
        paragraphs = shape.text_frame.paragraphs
        if paragraphs:
            first_paragraph = paragraphs[0]
            if first_paragraph.runs:
                first_paragraph.runs[0].text = new_text
            else:
                first_paragraph.add_run().text = new_text
        else:
            p = shape.text_frame.add_paragraph()
            p.add_run().text = new_text

    # ---------------- PPT保存 ----------------
    def save_ppt(self):
        if not self.ppt_path:
            return
        text = self.input_text.toPlainText()
        slides_text = text.split('--- Slide ')[1:]

        prs = Presentation(self.ppt_path)
        for i, slide in enumerate(prs.slides):
            if i >= len(slides_text):
                continue
            lines = slides_text[i].split('---\n', 1)[-1].strip()
            shapes_text = [s for s in slide.shapes if hasattr(s, "text") and s.has_text_frame]
            if shapes_text:
                self.replace_text_keep_font(shapes_text[0], lines)
            else:
                textbox = slide.shapes.add_textbox(
                    left=Inches(1), top=Inches(1), width=Inches(8), height=Inches(2)
                )
                self.replace_text_keep_font(textbox, lines)

        self.edited_ppt_path = self.ppt_path.replace(".pptx", "_edited.pptx")
        prs.save(self.edited_ppt_path)
        self.output_text.append(f"PPT出力完了: {self.edited_ppt_path}")

        # 編集後のスライドを表示
        if self.view_box.currentText() == "編集後PPT":
            self.convert_ppt_to_images(self.edited_ppt_path)
            self.show_slide()

    # ---------------- PPT → PNG変換 ----------------
    def convert_ppt_to_images(self, ppt_path):
        out_dir = os.path.join(os.path.dirname(ppt_path), "slides_tmp")
        os.makedirs(out_dir, exist_ok=True)
        subprocess.run([
            "soffice", "--headless", "--convert-to", "png",
            "--outdir", out_dir, ppt_path
        ], check=True)
        self.slides_pix = sorted([
            os.path.join(out_dir, f) for f in os.listdir(out_dir) if f.endswith(".png")
        ])
        self.slide_index = 0

    # ---------------- スライド表示 ----------------
    def show_slide(self):
        if not self.slides_pix:
            return
        pix = QPixmap(self.slides_pix[self.slide_index])
        self.slide_label.setPixmap(pix)

    # ---------------- 表示切替 ----------------
    def update_slide_view(self):
        if self.view_box.currentText() == "元のPPT" and self.ppt_path:
            self.convert_ppt_to_images(self.ppt_path)
        elif self.view_box.currentText() == "編集後PPT" and self.edited_ppt_path:
            self.convert_ppt_to_images(self.edited_ppt_path)
        self.show_slide()

# ---------------- アプリ起動 ----------------
if __name__ == "__main__":
    app = QApplication(sys.argv)
    try:
        with open("style.qss", "r", encoding="utf-8") as f:
            app.setStyleSheet(f.read())
    except FileNotFoundError:
        pass

    window = PPTApp()
    window.show()
    sys.exit(app.exec())
