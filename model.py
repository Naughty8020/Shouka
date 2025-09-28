from pptx import Presentation
from pptx.util import Inches

class PPTModel:
    def __init__(self, ppt_path):
        self.ppt_path = ppt_path
        self.edited_ppt_path = None
        self.presentation = Presentation(ppt_path)

    def extract_text(self):
        slides_text = []
        for i, slide in enumerate(self.presentation.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.has_text_frame:
                    texts.append(shape.text)
            slide_text = f"--- Slide {i + 1} ---\n" + "\n".join(texts)
            slides_text.append(slide_text)
        return "\n\n".join(slides_text)

    def save_with_texts(self, edited_text):
        slides_text = edited_text.split('--- Slide ')[1:]
        for i, slide in enumerate(self.presentation.slides):
            if i >= len(slides_text):
                continue
            text = slides_text[i].split('---\n', 1)[-1].strip()
            shapes = [s for s in slide.shapes if hasattr(s, "text") and s.has_text_frame]
            if shapes:
                self.replace_text_keep_font(shapes[0], text)
            else:
                textbox = slide.shapes.add_textbox(
                    Inches(1), Inches(1), Inches(8), Inches(2)
                )
                self.replace_text_keep_font(textbox, text)
        self.edited_ppt_path = self.ppt_path.replace(".pptx", "_edited.pptx")
        self.presentation.save(self.edited_ppt_path)
        return self.edited_ppt_path

    def replace_text_keep_font(self, shape, new_text):
        if not shape.has_text_frame:
            return
        paragraphs = shape.text_frame.paragraphs
        if paragraphs:
            first_paragraph = paragraphs[0]
            if first_paragraph.runs:
                first_paragraph.runs[0].text = new_text
            else:
                first_paragraph.add_run().text = new_text
